"""
generate_pptx.py — Genera la presentacion ejecutiva de SmarTrack como .pptx
Mantiene el mismo contenido y estructura del HTML, con diseno oscuro premium.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree

# ── PALETTE ──────────────────────────────────────────────
BG        = RGBColor(0x06, 0x08, 0x10)
SURFACE   = RGBColor(0x0d, 0x11, 0x17)
SURFACE2  = RGBColor(0x16, 0x1b, 0x22)
PRIMARY   = RGBColor(0x3b, 0x82, 0xf6)
ACCENT    = RGBColor(0x63, 0x66, 0xf1)
TEAL      = RGBColor(0x14, 0xb8, 0xa6)
GOLD      = RGBColor(0xf5, 0x9e, 0x0b)
SUCCESS   = RGBColor(0x10, 0xb9, 0x81)
DANGER    = RGBColor(0xef, 0x44, 0x44)
WHITE     = RGBColor(0xf0, 0xf6, 0xfc)
SOFT      = RGBColor(0x8b, 0x94, 0x9e)
MUTED     = RGBColor(0x48, 0x4f, 0x58)
BORDER    = RGBColor(0x1c, 0x21, 0x28)

# Slide size: widescreen 16:9
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# Blank layout
blank = prs.slide_layouts[6]

# ── HELPER FUNCTIONS ────────────────────────────────────

def add_bg(slide, color=BG):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0), radius=None):
    """Add a rounded rectangle shape."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()

    if radius:
        # Apply rounded corners via XML
        sp = shape._element
        prstGeom = sp.find(qn('p:spPr')).find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            else:
                for av in avLst.findall(qn('a:gd')):
                    avLst.remove(av)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    return shape


def txt(slide, text, left, top, width, height,
        font_size=Pt(14), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return txb


def txt_multi(slide, lines, left, top, width, height,
              font_size=Pt(12), bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, line_spacing=None):
    """Add multi-line text box (list of (text, color, bold, size) tuples)."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, color, bold, font_size)
        t, c, b, sz = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing

        run = p.add_run()
        run.text = t
        run.font.size = sz
        run.font.bold = b
        run.font.color.rgb = c
        run.font.name = 'Calibri'

    return txb


def gradient_rect(slide, left, top, width, height, color1, color2, radius=None):
    """Approximate gradient with semi-transparent overlay boxes."""
    b1 = box(slide, left, top, width, height, fill_color=color1, radius=radius)
    b2 = box(slide, left, top, width, height, fill_color=color2, radius=radius)
    sp = b2._element
    spPr = sp.find(qn('p:spPr'))
    fill = spPr.find(qn('a:solidFill'))
    if fill is not None:
        spPr.remove(fill)
    gradFill = etree.SubElement(spPr, qn('a:gradFill'))
    gradFill.set('rotWithShape', '1')
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

    gs1 = etree.SubElement(gsLst, qn('a:gs'))
    gs1.set('pos', '0')
    solidFill1 = etree.SubElement(gs1, qn('a:solidFill'))
    srgb1 = etree.SubElement(solidFill1, qn('a:srgbClr'))
    srgb1.set('val', f'{color1.red:02X}{color1.green:02X}{color1.blue:02X}')

    gs2 = etree.SubElement(gsLst, qn('a:gs'))
    gs2.set('pos', '100000')
    solidFill2 = etree.SubElement(gs2, qn('a:solidFill'))
    srgb2 = etree.SubElement(solidFill2, qn('a:srgbClr'))
    srgb2.set('val', f'{color2.red:02X}{color2.green:02X}{color2.blue:02X}')

    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', '5400000')  # 90 deg = top to bottom; 16200000 = left to right
    lin.set('scaled', '0')

    b2.line.fill.background()
    return b1


def card_block(slide, left, top, width, height, icon, title, desc,
               title_color=WHITE, bg_color=SURFACE2, border_color=BORDER,
               icon_size=Pt(24), title_size=Pt(14), desc_size=Pt(10.5)):
    """Draw a card with icon, title, description."""
    box(slide, left, top, width, height, fill_color=bg_color,
        line_color=border_color, line_width=Pt(0.5), radius=30000)

    # icon
    txt(slide, icon, left + Inches(0.25), top + Inches(0.2),
        Inches(0.6), Inches(0.5), font_size=icon_size, align=PP_ALIGN.LEFT)

    # title
    txt(slide, title, left + Inches(0.25), top + Inches(0.75),
        width - Inches(0.5), Inches(0.35),
        font_size=title_size, bold=True, color=title_color)

    # desc
    txt(slide, desc, left + Inches(0.25), top + Inches(1.15),
        width - Inches(0.5), height - Inches(1.25),
        font_size=desc_size, color=SOFT)


def eyebrow(slide, text, left, top, color=PRIMARY):
    txt(slide, text.upper(), left, top, Inches(6), Inches(0.3),
        font_size=Pt(9), bold=True, color=color)


def checklist_item(slide, text, left, top, width, color=SOFT, size=Pt(11)):
    """Draw a check mark + text."""
    txt(slide, '✓', left, top, Inches(0.3), Inches(0.35),
        font_size=size, bold=True, color=SUCCESS)
    txt(slide, text, left + Inches(0.3), top, width - Inches(0.3), Inches(0.45),
        font_size=size, color=color)


def slide_number_bg(slide, num):
    """Big faint slide number watermark."""
    txt(slide, str(num).zfill(2),
        W - Inches(2.5), Inches(0.2), Inches(2.4), Inches(1.5),
        font_size=Pt(90), bold=True, color=RGBColor(0x1a, 0x1f, 0x27),
        align=PP_ALIGN.RIGHT)


def stat_box(slide, left, top, w, h, number, label, num_color=PRIMARY):
    box(slide, left, top, w, h, fill_color=SURFACE2,
        line_color=BORDER, line_width=Pt(0.5), radius=20000)
    txt(slide, number, left, top + Inches(0.25), w, Inches(0.8),
        font_size=Pt(38), bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(slide, label, left + Inches(0.15), top + Inches(1.05),
        w - Inches(0.3), Inches(0.5), font_size=Pt(9), color=SOFT,
        align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 1 — HERO
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, BG)

# Background accent
box(s, Inches(-1), Inches(-1), Inches(8), Inches(6), fill_color=RGBColor(0x0a, 0x12, 0x25), radius=50000)

# Top glow bar
box(s, Inches(3), Inches(0.4), Inches(7), Pt(2), fill_color=PRIMARY)

# Logo mark
b = box(s, Inches(1.2), Inches(1.0), Inches(0.7), Inches(0.7), fill_color=PRIMARY, radius=20000)

# Title
txt(s, 'SmarTrack', Inches(1.2), Inches(1.9), Inches(10), Inches(1.4),
    font_size=Pt(64), bold=True, color=PRIMARY)
txt(s, 'Inventario Inteligente para Planta de Manufactura',
    Inches(1.2), Inches(3.1), Inches(10), Inches(0.9),
    font_size=Pt(22), bold=False, color=WHITE)
txt(s, 'Transformamos la forma en que su planta administra, visualiza y extrae\ncomponentes electrónicos — en tiempo real, desde cualquier dispositivo.',
    Inches(1.2), Inches(3.9), Inches(9.5), Inches(1.0),
    font_size=Pt(13), color=SOFT)

# Stats row
stats = [
    ('100%', 'Sin papel'),
    ('<5s', 'Actualización'),
    ('24/7', 'Disponible'),
    ('0', 'Errores por datos viejos'),
]
for i, (num, lbl) in enumerate(stats):
    x = Inches(1.2) + i * Inches(2.8)
    box(s, x, Inches(5.3), Inches(2.5), Inches(1.3), fill_color=SURFACE2,
        line_color=BORDER, line_width=Pt(0.5), radius=15000)
    colors = [PRIMARY, ACCENT, TEAL, GOLD]
    txt(s, num, x, Inches(5.35), Inches(2.5), Inches(0.65),
        font_size=Pt(30), bold=True, color=colors[i], align=PP_ALIGN.CENTER)
    txt(s, lbl, x, Inches(5.95), Inches(2.5), Inches(0.4),
        font_size=Pt(9), color=SOFT, align=PP_ALIGN.CENTER)

txt(s, 'SR', Inches(1.25), Inches(1.05), Inches(0.7), Inches(0.55),
    font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 2 — EL PROBLEMA
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, SURFACE)
slide_number_bg(s, 2)

eyebrow(s, '🔴  El Problema', Inches(0.6), Inches(0.45), DANGER)
txt(s, '¿Qué pasaba antes de SmarTrack?', Inches(0.6), Inches(0.75), Inches(12), Inches(0.9),
    font_size=Pt(34), bold=True, color=WHITE)
txt(s, 'Cada minuto que un operario busca un componente es productividad perdida.',
    Inches(0.6), Inches(1.55), Inches(11), Inches(0.45), font_size=Pt(13), color=SOFT)

problems = [
    ('🔍', 'Búsqueda Manual', 'El personal recorría físicamente los racks para verificar stock.\n5–20 minutos perdidos buscando un rollo.'),
    ('📋', 'Registros en Papel', 'Si alguien no actualizaba la hoja, el inventario quedaba\ndesactualizado. Datos que no reflejan la realidad.'),
    ('📞', 'Comunicación por Radio', 'Operador llama al almacén, el almacenista va a verificar.\nProceso lento con riesgo de errores y demoras en línea.'),
    ('❌', 'Sin Trazabilidad', '¿Quién sacó qué material y cuándo? Imposible saberlo.\nAuditorías complicadas, mermas difíciles de justificar.'),
]

card_positions = [
    (Inches(0.5), Inches(2.2)),
    (Inches(6.7), Inches(2.2)),
    (Inches(0.5), Inches(4.6)),
    (Inches(6.7), Inches(4.6)),
]
for (icon, title, desc), (cx, cy) in zip(problems, card_positions):
    card_block(s, cx, cy, Inches(6.0), Inches(2.1), icon, title, desc,
               bg_color=RGBColor(0x14, 0x0a, 0x0a),
               border_color=RGBColor(0x4a, 0x15, 0x18))

# Callout
box(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5),
    fill_color=RGBColor(0x13, 0x09, 0x09), line_color=DANGER, line_width=Pt(0.5), radius=10000)
txt(s, '🔴  Costo estimado: 30–90 minutos por turno en búsquedas manuales que podrían eliminarse por completo.',
    Inches(0.7), Inches(6.88), Inches(12), Inches(0.4), font_size=Pt(10.5), color=DANGER)

# ═══════════════════════════════════════════════════════
# SLIDE 3 — LA SOLUCIÓN
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, BG)
slide_number_bg(s, 3)

eyebrow(s, '✅  La Solución', Inches(0.6), Inches(0.45), PRIMARY)
txt(s, 'SmarTrack conecta hardware y operación', Inches(0.6), Inches(0.75), Inches(12), Inches(0.85),
    font_size=Pt(34), bold=True, color=WHITE)
txt(s, 'El sistema lee los racks físicos automáticamente cada 5 segundos, mantiene\nuna base de datos siempre al día y la expone en un panel web accesible desde cualquier PC.',
    Inches(0.6), Inches(1.55), Inches(11), Inches(0.7), font_size=Pt(12.5), color=SOFT)

# Flow diagram
flow_items = [
    ('🏗️', 'Hardware', 'SmartRack\n& Torres JUKI'),
    ('🔄', 'Sincronización', 'Cada 5 segundos\nautomática'),
    ('💾', 'Base de Datos', 'Inventario local\nsiempre al día'),
    ('🌐', 'Panel Web', 'Cualquier\ndispositivo'),
    ('👤', 'Operador', 'Solicita en\nsegundos'),
]
cols = [PRIMARY, TEAL, ACCENT, PRIMARY, TEAL]
fx = Inches(0.4)
for i, (icon, title, sub) in enumerate(flow_items):
    bcolor = RGBColor(0x12, 0x1d, 0x34) if i % 2 == 0 else RGBColor(0x0e, 0x1e, 0x1e)
    lcolor = cols[i]
    box(s, fx, Inches(2.55), Inches(2.2), Inches(1.75), fill_color=bcolor,
        line_color=lcolor, line_width=Pt(0.8), radius=15000)
    txt(s, icon, fx, Inches(2.6), Inches(2.2), Inches(0.5),
        font_size=Pt(20), align=PP_ALIGN.CENTER)
    txt(s, title, fx, Inches(3.1), Inches(2.2), Inches(0.4),
        font_size=Pt(11), bold=True, color=cols[i], align=PP_ALIGN.CENTER)
    txt(s, sub, fx, Inches(3.45), Inches(2.2), Inches(0.7),
        font_size=Pt(9), color=SOFT, align=PP_ALIGN.CENTER)
    fx += Inches(2.2)
    if i < 4:
        txt(s, '→', Inches(0.4) + i * Inches(2.2) + Inches(2.1), Inches(3.1), Inches(0.2), Inches(0.4),
            font_size=Pt(16), color=MUTED, align=PP_ALIGN.CENTER)

# 3 pillars
pillars = [
    ('📡', 'Polling Automático', 'El sistema consulta los racks físicos cada 5 segundos sin intervención humana.'),
    ('🔀', 'Dos Tipos de Almacén', 'Integra SmartRack y Torres JUKI en una sola vista unificada y coherente.'),
    ('⚡', 'Sin Instalación', 'Se accede desde el navegador. No instalar nada en las PCs de producción.'),
]
px = Inches(0.5)
for icon, title, desc in pillars:
    box(s, px, Inches(4.65), Inches(4.1), Inches(2.1), fill_color=SURFACE2,
        line_color=BORDER, line_width=Pt(0.5), radius=15000)
    txt(s, icon, px + Inches(0.15), Inches(4.75), Inches(0.5), Inches(0.5),
        font_size=Pt(22))
    txt(s, title, px + Inches(0.15), Inches(5.3), Inches(3.7), Inches(0.35),
        font_size=Pt(12), bold=True, color=WHITE)
    txt(s, desc, px + Inches(0.15), Inches(5.6), Inches(3.7), Inches(0.9),
        font_size=Pt(10), color=SOFT)
    px += Inches(4.4)

# ═══════════════════════════════════════════════════════
# SLIDE 4 — PANEL OPERADOR
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, SURFACE)
slide_number_bg(s, 4)

eyebrow(s, '🔍  Panel de Operador', Inches(0.6), Inches(0.45), TEAL)
txt(s, 'El operario encuentra cualquier\ncomponente en segundos', Inches(0.6), Inches(0.75), Inches(7), Inches(1.5),
    font_size=Pt(28), bold=True, color=WHITE)

items = [
    'Escanea o escribe el número de parte — el sistema localiza rack y celda exacta.',
    'Soporte de scanner de barras (formato Matdoc) — búsqueda automática al escanear.',
    'Autocompletado inteligente mientras escribes con partes disponibles en inventario.',
    'Selección de Línea para priorizar los racks correctos para esa línea de producción.',
    'Buffer multi-rollo — agrega varios componentes y extrae todos de un solo clic.',
    'Alerta visual en rojo si el componente está en un rack de otra línea.',
]
for i, item in enumerate(items):
    checklist_item(s, item, Inches(0.6), Inches(2.3) + i * Inches(0.62), Inches(6.5), size=Pt(11))

# Mock screen (right side)
box(s, Inches(7.7), Inches(1.0), Inches(5.3), Inches(6.1), fill_color=SURFACE,
    line_color=BORDER, line_width=Pt(0.8), radius=20000)

# mock bar
box(s, Inches(7.7), Inches(1.0), Inches(5.3), Inches(0.38), fill_color=SURFACE2, radius=20000)
for xi, dc in enumerate([DANGER, GOLD, SUCCESS]):
    dot = box(s, Inches(7.85) + xi * Inches(0.22), Inches(1.13), Inches(0.12), Inches(0.12),
              fill_color=dc, radius=50000)

txt(s, 'SmarTrack — Panel Operador', Inches(8.2), Inches(1.05), Inches(4), Inches(0.3),
    font_size=Pt(9), color=MUTED)

# search input mock
box(s, Inches(7.85), Inches(1.55), Inches(4.9), Inches(0.55), fill_color=BG,
    line_color=TEAL, line_width=Pt(1.2), radius=10000)
txt(s, '🔍  0402WGF10R0TCE', Inches(8.0), Inches(1.6), Inches(4.5), Inches(0.4),
    font_size=Pt(13), bold=True, color=TEAL)

# result found
box(s, Inches(7.85), Inches(2.25), Inches(4.9), Inches(1.0),
    fill_color=RGBColor(0x05, 0x18, 0x12),
    line_color=SUCCESS, line_width=Pt(0.8), radius=10000)
txt(s, '✓  ENCONTRADO EN LÍNEA', Inches(8.0), Inches(2.3), Inches(4.5), Inches(0.35),
    font_size=Pt(9.5), bold=True, color=SUCCESS)
txt(s, 'SMARTRACK 2  ·  Posición Left B/14', Inches(8.0), Inches(2.6), Inches(4.5), Inches(0.35),
    font_size=Pt(12), bold=True, color=WHITE)
txt(s, 'Qty: 4,200 uds  ·  Actualizado hace 3s', Inches(8.0), Inches(2.9), Inches(4.5), Inches(0.3),
    font_size=Pt(9), color=MUTED)

# warning
box(s, Inches(7.85), Inches(3.4), Inches(4.9), Inches(0.55),
    fill_color=RGBColor(0x18, 0x08, 0x08),
    line_color=DANGER, line_width=Pt(0.7), radius=10000)
txt(s, '⚠️  1 rollo adicional en SMARTRACK 5 (otra línea)',
    Inches(8.0), Inches(3.47), Inches(4.5), Inches(0.35), font_size=Pt(10), color=DANGER)

# buffer items
txt(s, 'Buffer de Extracción  (2 rollos)', Inches(7.85), Inches(4.1), Inches(4.9), Inches(0.35),
    font_size=Pt(10), bold=True, color=SOFT)
for bi, (pn, loc, qty) in enumerate([
    ('0402WGF10R0TCE', 'SR2 · Left B/14', '4,200'),
    ('RC0402FR-07100RL', 'SR1 · Right A/05', '12,000'),
]):
    by = Inches(4.5) + bi * Inches(0.65)
    box(s, Inches(7.85), by, Inches(4.9), Inches(0.55), fill_color=BG,
        line_color=BORDER, line_width=Pt(0.5), radius=8000)
    txt(s, pn, Inches(8.0), by + Inches(0.05), Inches(2.2), Inches(0.3),
        font_size=Pt(9.5), bold=True, color=PRIMARY)
    txt(s, loc, Inches(8.0), by + Inches(0.3), Inches(2.5), Inches(0.2),
        font_size=Pt(8.5), color=MUTED)
    txt(s, qty, Inches(11.7), by + Inches(0.15), Inches(0.9), Inches(0.25),
        font_size=Pt(9.5), bold=True, color=SUCCESS, align=PP_ALIGN.RIGHT)

# Extract button
box(s, Inches(7.85), Inches(5.85), Inches(4.9), Inches(0.7), fill_color=SUCCESS,
    radius=10000)
txt(s, '⚡  EXTRAER SELECCIONADOS', Inches(7.85), Inches(5.95), Inches(4.9), Inches(0.5),
    font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 5 — EXTRACCIÓN
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, BG)
slide_number_bg(s, 5)

eyebrow(s, '⚡  Extracción de Material', Inches(0.6), Inches(0.45), PRIMARY)
txt(s, 'Un clic para extraer — inmediato o programado', Inches(0.6), Inches(0.75), Inches(12), Inches(0.8),
    font_size=Pt(30), bold=True, color=WHITE)
txt(s, 'Elige cuándo se ejecuta la extracción: al instante o con anticipación.',
    Inches(0.6), Inches(1.5), Inches(11), Inches(0.4), font_size=Pt(13), color=SOFT)

ext_cards = [
    ('⚡', 'Extracción Inmediata', PRIMARY,
     'El sistema envía la señal al rack físico al instante. El rack libera los rollos y el operador solo va a recogerlos. Sin esperas, sin movimientos innecesarios.'),
    ('⏰', 'Extracción Programada', ACCENT,
     'Programa la extracción con anticipación: en 5, 10, 15, 30 minutos o 1 hora. El sistema la ejecuta automáticamente a la hora exacta que definiste.'),
    ('🏢', 'Pedidos a Torre JUKI', TEAL,
     'Si el componente está en la torre JUKI, genera un ticket de pedido al operador de la torre con nivel de urgencia configurable del 1 al 5.'),
    ('🔁', 'Anti-Duplicados', GOLD,
     'Si el internet falla y el usuario presiona dos veces, el sistema detecta la solicitud duplicada y retorna el resultado original sin ejecutar dos veces.'),
]

positions = [
    (Inches(0.5), Inches(2.15)),
    (Inches(6.65), Inches(2.15)),
    (Inches(0.5), Inches(4.7)),
    (Inches(6.65), Inches(4.7)),
]
for (icon, title, color, desc), (cx, cy) in zip(ext_cards, positions):
    bcolor = RGBColor(0x0d, 0x14, 0x24) if color == PRIMARY else \
             RGBColor(0x0f, 0x0f, 0x20) if color == ACCENT else \
             RGBColor(0x08, 0x18, 0x18) if color == TEAL else \
             RGBColor(0x18, 0x12, 0x05)
    card_block(s, cx, cy, Inches(6.0), Inches(2.3), icon, title, desc,
               title_color=color, bg_color=bcolor, border_color=color,
               icon_size=Pt(22), title_size=Pt(14), desc_size=Pt(10.5))

# Callout bottom
box(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.37),
    fill_color=RGBColor(0x0e, 0x14, 0x24), line_color=PRIMARY, line_width=Pt(0.5), radius=8000)
txt(s, '📊  Panel visual de programadas con cuenta regresiva en tiempo real y opción de cancelar antes de ejecutarse.',
    Inches(0.7), Inches(7.02), Inches(12), Inches(0.32), font_size=Pt(10), color=PRIMARY)

# ═══════════════════════════════════════════════════════
# SLIDE 6 — INVENTARIO EN VIVO
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, SURFACE)
slide_number_bg(s, 6)

eyebrow(s, '📦  Inventario en Vivo', Inches(0.6), Inches(0.45), ACCENT)
txt(s, 'Visibilidad total del stock,\nsiempre actualizada', Inches(0.6), Inches(0.75), Inches(6.8), Inches(1.4),
    font_size=Pt(30), bold=True, color=WHITE)

inv_items = [
    'Todos los rollos con código, número de parte, cantidad, posición exacta y hora de actualización.',
    'Buscador instantáneo por código de rollo o número de parte.',
    'Exporta todo el inventario a Excel/CSV con un solo clic.',
    'Semáforo Verde "En Vivo" para confirmar que los datos son actuales.',
    'Incluye inventario de Torres JUKI fusionado en la misma vista.',
    'Se refresca automáticamente cada 5 segundos sin recargar la página.',
]
for i, itm in enumerate(inv_items):
    checklist_item(s, itm, Inches(0.6), Inches(2.3) + i * Inches(0.62), Inches(6.5), size=Pt(11))

# Mock inventory table
box(s, Inches(7.5), Inches(1.0), Inches(5.5), Inches(6.3), fill_color=SURFACE,
    line_color=BORDER, line_width=Pt(0.8), radius=15000)

# header bar
box(s, Inches(7.5), Inches(1.0), Inches(5.5), Inches(0.45), fill_color=SURFACE2, radius=15000)
txt(s, 'Inventario DB — SmarTrack', Inches(7.65), Inches(1.06), Inches(3), Inches(0.3),
    font_size=Pt(9.5), bold=True, color=WHITE)
txt(s, '● En vivo', Inches(11.2), Inches(1.06), Inches(1.6), Inches(0.3),
    font_size=Pt(9), color=SUCCESS, align=PP_ALIGN.RIGHT)

# table header
for label, cx in [('Part Number', Inches(7.6)), ('Posicion', Inches(9.6)), ('Qty', Inches(11.8))]:
    txt(s, label, cx, Inches(1.55), Inches(2.0), Inches(0.3),
        font_size=Pt(8), bold=True, color=MUTED)

rows = [
    ('0402WGF10R0TCE',     'SR2 · Left B/14',   '4,200',  PRIMARY),
    ('RC0402FR-07100RL',   'SR1 · Right A/05',  '12,000', PRIMARY),
    ('GRM155R71C104KA88D', 'JUKI Torre 3',       '850',    GOLD),
    ('ERJ-2RKF100X',       'SR3 · Left C/22',   '2,750',  PRIMARY),
    ('CC0402KRX7R9BB104',  'SR2 · Right D/07',  '6,100',  PRIMARY),
    ('CRCW0402100KFKED',   'JUKI Torre 1',       '3,400',  GOLD),
    ('RC0402JR-07100RL',   'SR4 · Left A/11',   '900',    PRIMARY),
]
for ri, (pn, pos, qty, pc) in enumerate(rows):
    ry = Inches(1.9) + ri * Inches(0.58)
    if ri % 2 == 0:
        box(s, Inches(7.5), ry, Inches(5.5), Inches(0.57),
            fill_color=RGBColor(0x10, 0x15, 0x1c))
    txt(s, pn, Inches(7.6), ry + Inches(0.1), Inches(1.9), Inches(0.35),
        font_size=Pt(9), color=pc)
    txt(s, pos, Inches(9.55), ry + Inches(0.1), Inches(2.0), Inches(0.35),
        font_size=Pt(9), color=SOFT)
    txt(s, qty, Inches(11.7), ry + Inches(0.1), Inches(1.1), Inches(0.35),
        font_size=Pt(9), bold=True, color=SUCCESS if pc == PRIMARY else GOLD,
        align=PP_ALIGN.RIGHT)

# CSV button
box(s, Inches(7.6), Inches(6.7), Inches(2.2), Inches(0.45),
    fill_color=RGBColor(0x0e, 0x18, 0x2e), line_color=PRIMARY, line_width=Pt(0.7), radius=8000)
txt(s, '⬇  Exportar CSV', Inches(7.6), Inches(6.75), Inches(2.2), Inches(0.35),
    font_size=Pt(10), bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 7 — TRAZABILIDAD
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, BG)
slide_number_bg(s, 7)

eyebrow(s, '🔐  Control & Trazabilidad', Inches(0.6), Inches(0.45), GOLD)
txt(s, 'Saber quién sacó qué.\nSiempre.', Inches(0.6), Inches(0.75), Inches(12), Inches(1.3),
    font_size=Pt(36), bold=True, color=WHITE)
txt(s, 'Cada extracción queda registrada con nombre del operador, hora y resultado.',
    Inches(0.6), Inches(1.95), Inches(11), Inches(0.4), font_size=Pt(13), color=SOFT)

traz_cards = [
    ('🔐', 'Acceso con Número de Empleado', RGBColor(0x1a, 0x12, 0x04), GOLD,
     'Cada operador usa su número de empleado. Queda registrado quién solicitó cada movimiento de material.'),
    ('📜', 'Auditoría Completa', RGBColor(0x10, 0x10, 0x22), ACCENT,
     'El panel administrativo muestra el historial de extracciones: operador, resultado, ID del trabajo y hora exacta.'),
    ('🏢', 'Panel Operador JUKI', RGBColor(0x08, 0x18, 0x18), TEAL,
     'Operadores de la torre JUKI tienen su propio panel donde ven pedidos entrantes, urgencia y los marcan como completados.'),
    ('🔑', 'Gestión de Accesos', RGBColor(0x0e, 0x18, 0x2e), PRIMARY,
     'Registra empleados, asigna o revoca accesos en segundos. Si alguien sale de la empresa, su acceso se deshabilita de inmediato.'),
]
tpos = [
    (Inches(0.5), Inches(2.6)),
    (Inches(6.65), Inches(2.6)),
    (Inches(0.5), Inches(5.1)),
    (Inches(6.65), Inches(5.1)),
]
for (icon, title, bg, bc, desc), (cx, cy) in zip(traz_cards, tpos):
    card_block(s, cx, cy, Inches(6.0), Inches(2.25), icon, title, desc,
               title_color=bc, bg_color=bg, border_color=bc)

# ═══════════════════════════════════════════════════════
# SLIDE 8 — LISTO PARA PRODUCCIÓN
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, SURFACE)
slide_number_bg(s, 8)

eyebrow(s, '🛡️  Listo para Producción', Inches(0.6), Inches(0.45), ACCENT)
txt(s, 'Un sistema construido para\nresistir la planta', Inches(0.6), Inches(0.75), Inches(12), Inches(1.3),
    font_size=Pt(34), bold=True, color=WHITE)
txt(s, 'No es un prototipo. SmarTrack fue diseñado y probado para operar en entornos industriales reales.',
    Inches(0.6), Inches(1.95), Inches(12), Inches(0.4), font_size=Pt(13), color=SOFT)

# Stats
kpis = [
    ('85%', 'Cobertura de pruebas automáticas', PRIMARY),
    ('130', 'Tests unitarios y de integración', SUCCESS),
    ('24/7', 'Funcionamiento continuo', TEAL),
    ('.exe', 'Sin instalación compleja', GOLD),
]
for i, (num, lbl, c) in enumerate(kpis):
    stat_box(s, Inches(0.5) + i * Inches(3.1), Inches(2.65), Inches(2.9), Inches(1.5), num, lbl, c)

# Two columns
left_items = [
    ('Circuit Breaker:', 'Si el SmartRack no responde, el sistema no colapsa. Reanuda automáticamente.'),
    ('Reconexión automática:', 'Al servidor SmartRack si el token expira durante operación.'),
    ('Modo lectura:', 'Si hay problemas de escritura, jamás pierde visibilidad del inventario.'),
    ('Limpieza automática:', 'Registros viejos se eliminan solos para mantener el rendimiento.'),
]
right_items = [
    ('Safe Mode:', 'Se niega a iniciar si detecta contraseñas débiles por defecto.'),
    ('Capas de acceso:', 'Operadores usan número de empleado; admins tienen credenciales separadas.'),
    ('API Keys revocables:', 'Si un empleado sale, su acceso se deshabilita en segundos.'),
    ('Logs estructurados:', 'Con ID de correlación por cada solicitud para diagnóstico rápido.'),
]

for col_idx, (col_items, col_title, col_x) in enumerate([
    (left_items, 'RESILIENCIA Y RECUPERACIÓN', Inches(0.6)),
    (right_items, 'SEGURIDAD DE DESPLIEGUE', Inches(6.8)),
]):
    txt(s, col_title, col_x, Inches(4.45), Inches(6), Inches(0.3),
        font_size=Pt(8.5), bold=True, color=MUTED)
    for i, (label, desc) in enumerate(col_items):
        ty = Inches(4.85) + i * Inches(0.64)
        txt(s, f'✓  {label}', col_x, ty, Inches(6), Inches(0.3),
            font_size=Pt(10.5), bold=True, color=SUCCESS)
        txt(s, desc, col_x + Inches(0.2), ty + Inches(0.28), Inches(5.8), Inches(0.3),
            font_size=Pt(10), color=SOFT)

# ═══════════════════════════════════════════════════════
# SLIDE 9 — DESPLIEGUE SIMPLE
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, BG)
slide_number_bg(s, 9)

eyebrow(s, '🚀  Implementación Simple', Inches(0.6), Inches(0.45), TEAL)
txt(s, 'Funciona desde el primer día\nsin complicaciones de IT', Inches(0.6), Inches(0.75), Inches(12), Inches(1.3),
    font_size=Pt(34), bold=True, color=WHITE)
txt(s, 'No hay nube. No hay licencias anuales. No hay servidores complejos. Software que vive en tu planta.',
    Inches(0.6), Inches(1.95), Inches(12), Inches(0.4), font_size=Pt(13), color=SOFT)

steps = [
    ('1', 'Configura el .env', PRIMARY, 'Editas un archivo de texto con la IP del SmartRack y las credenciales. Toma 2 minutos. No requiere conocimientos técnicos.'),
    ('2', 'Ejecutas el .exe', ACCENT, 'Doble clic en SmarTrack.exe. El sistema arranca, conecta al hardware y sincroniza. Sin instaladores, sin permisos especiales.'),
    ('3', 'Abres el navegador', TEAL, 'Cualquier PC, tablet o celular en la misma red accede escribiendo la IP en el navegador. Listo para usar desde el día 1.'),
]
for i, (num, title, color, desc) in enumerate(steps):
    cx = Inches(0.5) + i * Inches(4.3)
    box(s, cx, Inches(2.7), Inches(4.0), Inches(2.5),
        fill_color=RGBColor(0x0d, 0x14, 0x1e), line_color=color, line_width=Pt(0.8), radius=20000)
    box(s, cx + Inches(0.2), Inches(2.9), Inches(0.65), Inches(0.65), fill_color=color, radius=15000)
    txt(s, num, cx + Inches(0.2), Inches(2.92), Inches(0.65), Inches(0.55),
        font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, cx + Inches(0.2), Inches(3.65), Inches(3.5), Inches(0.4),
        font_size=Pt(14), bold=True, color=color)
    txt(s, desc, cx + Inches(0.2), Inches(4.05), Inches(3.6), Inches(1.0),
        font_size=Pt(10.5), color=SOFT)

# On-premise callout
box(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.3),
    fill_color=RGBColor(0x08, 0x15, 0x1e), line_color=PRIMARY, line_width=Pt(0.8), radius=15000)
txt(s, '🔒', Inches(0.75), Inches(5.7), Inches(0.5), Inches(0.95), font_size=Pt(26))
txt(s, 'Completamente On-Premise', Inches(1.3), Inches(5.6), Inches(10), Inches(0.4),
    font_size=Pt(14), bold=True, color=WHITE)
txt(s,
    'Todos los datos se guardan en tu propia computadora, en tu red. Sin datos en la nube. Sin dependencia de internet. Sin suscripciones mensuales. Lo que compras es tuyo completo.',
    Inches(1.3), Inches(6.0), Inches(11.0), Inches(0.75), font_size=Pt(11), color=SOFT)

# ═══════════════════════════════════════════════════════
# SLIDE 10 — ROI
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, SURFACE)
slide_number_bg(s, 10)

eyebrow(s, '💰  Retorno de Inversión', Inches(0.6), Inches(0.32), GOLD)
txt(s, 'Lo que cambia con SmarTrack', Inches(0.6), Inches(0.6), Inches(12), Inches(0.75),
    font_size=Pt(32), bold=True, color=WHITE)

# Table
headers = ['Situación', '❌  Sin SmarTrack', '✅  Con SmarTrack']
header_colors = [MUTED, DANGER, SUCCESS]
col_widths = [Inches(4.0), Inches(4.0), Inches(4.7)]
col_x = [Inches(0.5), Inches(4.5), Inches(8.5)]

# header row
box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.45), fill_color=SURFACE2,
    line_color=BORDER, line_width=Pt(0.3), radius=5000)
for hd, hc, cx, cw in zip(headers, header_colors, col_x, col_widths):
    txt(s, hd, cx + Inches(0.1), Inches(1.55), cw, Inches(0.35),
        font_size=Pt(9), bold=True, color=hc)

table_rows = [
    ('Localizar un componente', '5–20 minutos', '< 10 segundos'),
    ('Exactitud del inventario', 'Depende de quién anotó', 'Sincronizado en tiempo real'),
    ('Trazabilidad', 'Bitácora manual o ninguna', '100% automática con nombre'),
    ('Doble extracción', 'Riesgo alto', 'Eliminada por el sistema'),
    ('Solicitudes JUKI urgentes', 'Por radio o teléfono', 'Ticket digital con prioridad'),
    ('Extracción programada', 'No existe', 'Automática con anticipación'),
    ('Reportes de inventario', 'Horas de trabajo manual', 'Exportar CSV en 1 clic'),
    ('Dependencia de una persona', 'Alta — "el que sabe dónde está"', 'Cualquier operador puede buscar'),
]
for ri, (sit, bad, good) in enumerate(table_rows):
    ry = Inches(2.0) + ri * Inches(0.58)
    row_bg = RGBColor(0x10, 0x14, 0x1c) if ri % 2 == 0 else SURFACE
    box(s, Inches(0.5), ry, Inches(12.3), Inches(0.56), fill_color=row_bg)
    txt(s, sit, Inches(0.6), ry + Inches(0.1), Inches(3.8), Inches(0.35),
        font_size=Pt(10), color=SOFT)
    txt(s, bad, Inches(4.6), ry + Inches(0.1), Inches(3.8), Inches(0.35),
        font_size=Pt(10), color=DANGER)
    txt(s, good, Inches(8.6), ry + Inches(0.1), Inches(4.5), Inches(0.35),
        font_size=Pt(10), bold=True, color=SUCCESS)

# Bottom stats
kpis2 = [('-90%', 'Tiempo en localizar material', SUCCESS),
          ('+100%', 'Visibilidad del inventario', PRIMARY),
          ('0', 'Dependencias externas de software', GOLD)]
for i, (n, l, c) in enumerate(kpis2):
    stat_box(s, Inches(0.5) + i * Inches(4.1), Inches(6.75), Inches(3.9), Inches(0.6), n, '', c)
    txt(s, n, Inches(0.5) + i * Inches(4.1), Inches(6.77), Inches(3.9), Inches(0.45),
        font_size=Pt(20), bold=True, color=c, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 11 — RESUMEN EJECUTIVO
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, BG)
slide_number_bg(s, 11)

eyebrow(s, '📋  Resumen', Inches(0.6), Inches(0.4), ACCENT)
txt(s, 'Lo que obtiene con este desarrollo', Inches(0.6), Inches(0.7), Inches(12), Inches(0.8),
    font_size=Pt(32), bold=True, color=WHITE)

# Two columns
left_op = [
    'Encuentra cualquier componente en segundos desde su estación.',
    'Escanea con lector de barras para búsqueda instantánea.',
    'Solicita extracciones inmediatas o programadas con un clic.',
    'Ve extracciones pendientes con cuenta regresiva en tiempo real.',
    'Cierre de sesión seguro al final de cada turno.',
]
right_admin = [
    'Registra y gestiona el acceso de empleados desde la web.',
    'Revisa quién extrajo qué material y cuándo en la auditoría.',
    'Configura líneas de producción y sus racks asignados.',
    'Exporta inventario completo a Excel con un solo clic.',
    'Sin administración de servidores ni mantenimiento complejo.',
]

txt(s, 'Para el Operador', Inches(0.6), Inches(1.65), Inches(6.0), Inches(0.35),
    font_size=Pt(10), bold=True, color=TEAL)
for i, itm in enumerate(left_op):
    checklist_item(s, itm, Inches(0.6), Inches(2.1) + i * Inches(0.6), Inches(6.0), size=Pt(11))

txt(s, 'Para el Administrador', Inches(6.8), Inches(1.65), Inches(6.0), Inches(0.35),
    font_size=Pt(10), bold=True, color=PRIMARY)
for i, itm in enumerate(right_admin):
    checklist_item(s, itm, Inches(6.8), Inches(2.1) + i * Inches(0.6), Inches(6.0), size=Pt(11))

# Quote
box(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.2),
    fill_color=RGBColor(0x0b, 0x12, 0x22), line_color=PRIMARY, line_width=Pt(2), radius=15000)
box(s, Inches(0.5), Inches(5.2), Inches(0.08), Inches(1.2), fill_color=PRIMARY)
txt(s,
    '"Transformamos el inventario de algo que se manejaba con papel y radio, a un sistema digital,\ntrazable y automatizado — que cualquier operador puede usar desde el primer día."',
    Inches(0.75), Inches(5.3), Inches(12.0), Inches(1.0),
    font_size=Pt(12), italic=True, color=WHITE)

# Badges
badges = ['🏭 Manufactura electrónica', '📡 SmartRack & JUKI', '📊 Auditoría completa',
          '✅ 130 tests · Producción', '💾 On-Premise · Sin nube']
bx = Inches(0.5)
for b in badges:
    bw = Inches(2.3)
    box(s, bx, Inches(6.65), bw, Inches(0.45),
        fill_color=RGBColor(0x10, 0x18, 0x28), line_color=BORDER, line_width=Pt(0.5), radius=50000)
    txt(s, b, bx, Inches(6.68), bw, Inches(0.35), font_size=Pt(8.5), color=PRIMARY, align=PP_ALIGN.CENTER)
    bx += bw + Inches(0.08)

# ═══════════════════════════════════════════════════════
# SLIDE 12 — CIERRE
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, SURFACE)

# Center glow
box(s, Inches(4.5), Inches(1.5), Inches(4.3), Inches(4.5),
    fill_color=RGBColor(0x0a, 0x10, 0x28), radius=80000)

# Logo
box(s, Inches(5.8), Inches(1.4), Inches(1.7), Inches(1.1), fill_color=PRIMARY, radius=25000)
txt(s, 'SR', Inches(5.8), Inches(1.48), Inches(1.7), Inches(0.85),
    font_size=Pt(32), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title
txt(s, 'SmarTrack', Inches(2), Inches(2.7), Inches(9.3), Inches(1.0),
    font_size=Pt(54), bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
txt(s, 'Sistema de Inventario Inteligente', Inches(2), Inches(3.6), Inches(9.3), Inches(0.55),
    font_size=Pt(18), color=SOFT, align=PP_ALIGN.CENTER)
txt(s, 'Un sistema hecho a la medida de su operación, construido con tecnología\nmoderna, pensado para la realidad de su planta manufacturera.',
    Inches(2), Inches(4.2), Inches(9.3), Inches(0.9),
    font_size=Pt(12), color=SOFT, align=PP_ALIGN.CENTER)

# Ficha técnica
ficha = [
    ('Tecnología', 'FastAPI · SQLite · Python'),
    ('Acceso', 'Web · Cualquier navegador'),
    ('Hospedaje', 'On-Premise · Red local'),
    ('Calidad', '85% cobertura · 130 tests'),
]
fw = Inches(3.0)
fx = Inches(0.5)
fy = Inches(5.4)
for label, val in ficha:
    box(s, fx, fy, fw, Inches(0.85), fill_color=SURFACE2,
        line_color=BORDER, line_width=Pt(0.5), radius=10000)
    txt(s, label.upper(), fx + Inches(0.15), fy + Inches(0.07), fw - Inches(0.3), Inches(0.25),
        font_size=Pt(7.5), color=MUTED)
    txt(s, val, fx + Inches(0.15), fy + Inches(0.4), fw - Inches(0.3), Inches(0.35),
        font_size=Pt(11), bold=True, color=WHITE)
    fx += fw + Inches(0.1)

txt(s, '🔒 Todos los datos se almacenan en su infraestructura local. Sin dependencias externas.',
    Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.45),
    font_size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)

# ── SAVE ────────────────────────────────────────────────
out = r'c:\Proyectos\KimballInvertario\SmarTrack_Presentacion.pptx'
prs.save(out)
print(f'OK  Guardado: {out}')
